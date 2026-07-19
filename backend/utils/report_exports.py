"""
Report export engine.

CSV is always available (stdlib). Excel/PowerPoint/PDF are optional and guarded — if the
engine is not installed the caller gets a clear ExportUnavailable so the API can return 501
and the UI can disable that button (via the capabilities probe).
"""

import csv
import io

from utils.report_render import build_html

MIMES = {
    "html": "text/html",
    "csv": "text/csv",
    "xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "pdf": "application/pdf",
}


class ExportUnavailable(Exception):
    pass


# ---- capability probe (used by the router's /capabilities endpoint) ----
def available_formats():
    formats = {"html": True, "csv": True, "xlsx": False, "pptx": False, "pdf": False}
    try:
        import openpyxl  # noqa: F401
        formats["xlsx"] = True
    except Exception:
        pass
    try:
        import pptx  # noqa: F401
        formats["pptx"] = True
    except Exception:
        pass
    try:
        import weasyprint  # noqa: F401
        formats["pdf"] = True
    except Exception:
        pass
    return formats


# ---- flatten a definition into tabular rows ----
def _flatten(definition):
    rows = []  # (section, field, value)
    blocks = (definition or {}).get("blocks") or []
    for i, b in enumerate(blocks, start=1):
        kind = b.get("kind")
        if kind == "question":
            title = f"{i}. {b.get('q', '')}"
            if b.get("why"):
                rows.append((title, "Why", b["why"]))
            if b.get("metric"):
                rows.append((title, "Metric", b["metric"]))
            if b.get("decision"):
                rows.append((title, "Decision", b["decision"]))
        elif kind == "text":
            rows.append((f"{i}. {b.get('heading', 'Narrative')}", "Body", b.get("body", "")))
        elif kind == "kpis" or kind in ("live-exec", "live-kling", "live-cost"):
            title = f"{i}. {b.get('heading') or kind}"
            items = b.get("items") or b.get("snapshotItems") or []
            for it in items:
                rows.append((title, it.get("label", ""), it.get("value", "")))
        elif kind == "table":
            title = f"{i}. {b.get('title', 'Table')}"
            cols = b.get("columns") or []
            for r in b.get("rows") or []:
                rows.append((title, cols[0] if cols else "", " | ".join(str(c) for c in r)))
    return rows


def _report_title(definition):
    return ((definition or {}).get("branding") or {}).get("title") or "report"


# ---- exporters ----
def _csv_bytes(definition):
    buf = io.StringIO()
    w = csv.writer(buf)
    w.writerow(["Section", "Field", "Value"])
    for section, field, value in _flatten(definition):
        w.writerow([section, field, value])
    return buf.getvalue().encode("utf-8-sig")


def _xlsx_bytes(definition):
    try:
        from openpyxl import Workbook
        from openpyxl.styles import Font, PatternFill
    except Exception as exc:
        raise ExportUnavailable("Excel export requires the 'openpyxl' package.") from exc

    branding = (definition or {}).get("branding") or {}
    navy = (branding.get("navy") or "#101f3f").lstrip("#")
    wb = Workbook()
    ws = wb.active
    ws.title = "Report"
    ws.append([_report_title(definition)])
    ws["A1"].font = Font(size=14, bold=True)
    ws.append([])
    header = ["Section", "Field", "Value"]
    ws.append(header)
    hdr_fill = PatternFill("solid", fgColor=navy)
    for col in range(1, 4):
        cell = ws.cell(row=3, column=col)
        cell.font = Font(bold=True, color="FFFFFF")
        cell.fill = hdr_fill
    for section, field, value in _flatten(definition):
        ws.append([section, field, value])
    ws.column_dimensions["A"].width = 46
    ws.column_dimensions["B"].width = 18
    ws.column_dimensions["C"].width = 60
    out = io.BytesIO()
    wb.save(out)
    return out.getvalue()


def _pptx_bytes(definition):
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
    except Exception as exc:
        raise ExportUnavailable("PowerPoint export requires the 'python-pptx' package.") from exc

    branding = (definition or {}).get("branding") or {}
    navy_hex = (branding.get("navy") or "#101f3f").lstrip("#")
    gold_hex = (branding.get("gold") or "#c99a2e").lstrip("#")
    navy = RGBColor.from_string(navy_hex)
    gold = RGBColor.from_string(gold_hex)

    prs = Presentation()
    blank = prs.slide_layouts[6]

    # Title slide
    s = prs.slides.add_slide(blank)
    bg = s.background
    bg.fill.solid()
    bg.fill.fore_color.rgb = navy
    tb = s.shapes.add_textbox(Inches(0.7), Inches(2.2), Inches(9), Inches(2))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = branding.get("title") or "Corporate Report"
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = RGBColor.from_string("FFFFFF")
    p2 = tf.add_paragraph()
    p2.text = branding.get("subtitle") or "AI Intelligence Report"
    p2.font.size = Pt(16)
    p2.font.color.rgb = gold

    for i, b in enumerate(( definition or {}).get("blocks") or [], start=1):
        s = prs.slides.add_slide(blank)
        title = b.get("q") or b.get("heading") or b.get("title") or b.get("kind")
        tb = s.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(9), Inches(1))
        tp = tb.text_frame.paragraphs[0]
        tp.text = f"{i}. {title}"
        tp.font.size = Pt(24)
        tp.font.bold = True
        tp.font.color.rgb = navy
        body = s.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(8.6), Inches(4.5)).text_frame
        body.word_wrap = True
        lines = []
        if b.get("why"):
            lines.append(f"Why: {b['why']}")
        if b.get("metric"):
            lines.append(f"Metric: {b['metric']}")
        if b.get("decision"):
            lines.append(f"Decision: {b['decision']}")
        if b.get("body"):
            lines.append(b["body"])
        for it in (b.get("items") or b.get("snapshotItems") or []):
            lines.append(f"{it.get('label', '')}: {it.get('value', '')}")
        for r in (b.get("rows") or []):
            lines.append(" | ".join(str(c) for c in r))
        first = True
        for ln in lines or ["(no detail)"]:
            para = body.paragraphs[0] if first else body.add_paragraph()
            para.text = ln
            para.font.size = Pt(14)
            first = False

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def _pdf_bytes(definition, html_snapshot=None):
    try:
        from weasyprint import HTML
    except Exception as exc:
        raise ExportUnavailable(
            "Server-side PDF requires 'weasyprint'. Use the browser 'Print / Save as PDF' in the Report Builder, or install weasyprint on the server."
        ) from exc
    html = html_snapshot or build_html(definition)
    return HTML(string=html).write_pdf()


def export_report(definition, fmt, html_snapshot=None):
    """Return (bytes, mimetype, extension). Raises ExportUnavailable if engine missing."""
    fmt = (fmt or "").lower()
    if fmt == "html":
        return (html_snapshot or build_html(definition)).encode("utf-8"), MIMES["html"], "html"
    if fmt == "csv":
        return _csv_bytes(definition), MIMES["csv"], "csv"
    if fmt == "xlsx":
        return _xlsx_bytes(definition), MIMES["xlsx"], "xlsx"
    if fmt == "pptx":
        return _pptx_bytes(definition), MIMES["pptx"], "pptx"
    if fmt == "pdf":
        return _pdf_bytes(definition, html_snapshot), MIMES["pdf"], "pdf"
    raise ExportUnavailable(f"Unknown export format: {fmt}")
